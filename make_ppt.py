from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]

# Colors
DN=RGBColor(0x0C,0x1A,0x2E); TL=RGBColor(0x00,0xB4,0xD8); OR=RGBColor(0xFF,0x6B,0x35)
GR=RGBColor(0x06,0xD6,0xA0); WH=RGBColor(0xFF,0xFF,0xFF); LG=RGBColor(0xF0,0xF4,0xF8)
YL=RGBColor(0xF5,0x9E,0x0B); SL=RGBColor(0x94,0xA3,0xB8); RD=RGBColor(0xEF,0x44,0x44)
NV=RGBColor(0x0F,0x27,0x44); DT=RGBColor(0x1E,0x29,0x3B); PU=RGBColor(0x7C,0x3A,0xED)

def S(prs): return prs.slides.add_slide(BL)
def BG(sl,c):
    f=sl.background.fill; f.solid(); f.fore_color.rgb=c
def R(sl,l,t,w,h,fc=None,lc=None,lw=Pt(1)):
    sh=sl.shapes.add_shape(1,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.color.rgb=lc; sh.line.width=lw
    else: sh.line.fill.background()
    return sh
def T(sl,tx,l,t,w,h,sz,co,bo=False,al=PP_ALIGN.LEFT,it=False,fn="Calibri"):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=al; r=p.add_run()
    r.text=tx; r.font.size=Pt(sz); r.font.color.rgb=co
    r.font.bold=bo; r.font.italic=it; r.font.name=fn
    return tb
def ML(sl,lines,l,t,w,h,sz,co,bo=False,al=PP_ALIGN.LEFT,sp=Pt(4)):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=al; p.space_before=sp; r=p.add_run()
        r.text=ln; r.font.size=Pt(sz); r.font.color.rgb=co; r.font.bold=bo; r.font.name="Calibri"
    return tb
def HDR(sl,title,sub=None,dark=True):
    bg_c=DN if dark else LG; BG(sl,bg_c)
    R(sl,0,0,Inches(13.33),Inches(0.08),fc=TL)
    R(sl,0,Inches(7.42),Inches(13.33),Inches(0.08),fc=OR)
    tc=WH if dark else DT; sc=SL if dark else RGBColor(0x64,0x74,0x8B)
    T(sl,title,Inches(0.5),Inches(0.18),Inches(12),Inches(0.6),28,tc,True,PP_ALIGN.LEFT)
    if sub: T(sl,sub,Inches(0.5),Inches(0.72),Inches(12),Inches(0.35),12,sc,False,PP_ALIGN.LEFT)
    R(sl,Inches(0.5),Inches(1.0),Inches(12.33),Inches(0.03),fc=TL if dark else OR)

print("helpers loaded")
